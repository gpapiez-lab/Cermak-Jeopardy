#!/usr/bin/env python3
"""
PPTX → JSON (CONTENT ONLY)
- Reads a Jeopardy authoring deck with 2 rounds + Final.
- Extracts: category, question, answer, optional images/audio.
- NO game logic: no DD rules, no values, no scoring, no wagering.

Authoring conventions (recommended):
- Section header slide titles: "ROUND 1" and "ROUND 2"
- Final slide title: "FINAL JEOPARDY"
- Clue slides:
    - Slide TITLE = category
    - Slide BODY text = question (any non-title text on the slide)
    - Speaker Notes = answer
- Media:
    - Any picture shapes on slide become extracted images
    - Any embedded audio referenced by the slide is extracted

Usage:
  python3 ppt_to_json_content_only.py input.pptx --out game-data.json --assets assets

Outputs:
  game-data.json
  assets/images/...
  assets/audio/...
"""

import argparse
import json
import os
import re
import zipfile
import xml.etree.ElementTree as ET
from collections import OrderedDict

from pptx import Presentation

AUDIO_EXTS = {".mp3", ".m4a", ".wav", ".aac", ".ogg"}

def ensure_dir(p: str) -> None:
    os.makedirs(p, exist_ok=True)

def get_slide_title(slide) -> str:
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text_frame.text.strip()
    except Exception:
        pass
    return ""

def get_body_text(slide) -> str:
    """
    Return combined text from all non-title text frames.
    """
    parts = []
    title_shape = None
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if title_shape is not None and shape == title_shape:
            continue

        text = (shape.text_frame.text or "").strip()
        if text:
            parts.append(text)

    # Join multiple text boxes if present
    return "\n".join(parts).strip()

def get_notes_text(slide) -> str:
    """
    Speaker notes as answer.
    """
    try:
        notes = slide.notes_slide.notes_text_frame.text
        return (notes or "").strip()
    except Exception:
        return ""

def extract_images_from_slide(slide, images_dir: str, slide_idx_1based: int):
    """
    Extract picture shapes via python-pptx.
    Returns list of filenames (relative, inside images_dir).
    Note: if PowerPoint stored an image as EMF/WMF, it will export as .emf/.wmf.
    """
    extracted = []
    img_num = 0

    for shape in slide.shapes:
        # 13 == MSO_SHAPE_TYPE.PICTURE (avoid importing enums)
        if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
            img_num += 1
            img = shape.image
            ext = (img.ext or "png").lower()

            fname = f"slide{slide_idx_1based:03d}_img{img_num:02d}.{ext}"
            out_path = os.path.join(images_dir, fname)
            with open(out_path, "wb") as f:
                f.write(img.blob)

            extracted.append(fname)

    return extracted

def _rels_path_for_slide(slide_idx_1based: int) -> str:
    return f"ppt/slides/_rels/slide{slide_idx_1based}.xml.rels"

def _slide_audio_targets(pptx_path: str, slide_idx_1based: int):
    """
    Parse slide relationship XML to find media targets that look like audio.
    Returns list of internal zip paths like "ppt/media/media1.mp3"
    """
    rels_path = _rels_path_for_slide(slide_idx_1based)

    with zipfile.ZipFile(pptx_path, "r") as z:
        if rels_path not in z.namelist():
            return []
        xml_bytes = z.read(rels_path)

    root = ET.fromstring(xml_bytes)

    targets = []
    for rel in root:
        if not rel.tag.lower().endswith("relationship"):
            continue
        target = rel.attrib.get("Target", "") or ""
        if "media/" not in target:
            continue

        # Normalize "../media/media1.m4a" → "ppt/media/media1.m4a"
        norm = target.replace("\\", "/")
        norm = norm.replace("..", "").lstrip("/")
        if not norm.startswith("ppt/"):
            norm = "ppt/" + norm

        _, ext = os.path.splitext(norm.lower())
        if ext in AUDIO_EXTS:
            targets.append(norm)

    return targets

def extract_audio_from_slide(pptx_path: str, audio_dir: str, slide_idx_1based: int):
    """
    Extract any audio files referenced by this slide.
    Returns list of filenames saved in audio_dir.
    """
    targets = _slide_audio_targets(pptx_path, slide_idx_1based)
    if not targets:
        return []

    saved = []
    with zipfile.ZipFile(pptx_path, "r") as z:
        for t in targets:
            if t not in z.namelist():
                continue
            data = z.read(t)
            base = os.path.basename(t)  # e.g. media1.m4a
            out_name = f"slide{slide_idx_1based:03d}_{base}"
            out_path = os.path.join(audio_dir, out_name)
            with open(out_path, "wb") as f:
                f.write(data)
            saved.append(out_name)

    return saved

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", help="Input PPTX")
    ap.add_argument("--out", default="game-data.json", help="Output JSON path")
    ap.add_argument("--assets", default="assets", help="Assets folder (images/audio)")
    ap.add_argument("--round1_title", default="ROUND 1", help="Round 1 header slide title")
    ap.add_argument("--round2_title", default="ROUND 2", help="Round 2 header slide title")
    ap.add_argument("--final_title", default="FINAL JEOPARDY", help="Final slide title")
    args = ap.parse_args()

    pptx_path = args.pptx
    out_json = args.out
    assets_root = args.assets

    images_dir = os.path.join(assets_root, "images")
    audio_dir  = os.path.join(assets_root, "audio")
    ensure_dir(images_dir)
    ensure_dir(audio_dir)

    prs = Presentation(pptx_path)

    r1_key = args.round1_title.strip().upper()
    r2_key = args.round2_title.strip().upper()
    f_key  = args.final_title.strip().upper()

    current_round = None  # "round1" or "round2"
    # Use OrderedDict to preserve category order as encountered
    rounds = {
        "round1": OrderedDict(),  # category -> list of clues
        "round2": OrderedDict(),
    }
    final = None

    for idx, slide in enumerate(prs.slides, start=1):
        title = get_slide_title(slide).strip()
        upper = title.upper()

        if upper == r1_key:
            current_round = "round1"
            continue
        if upper == r2_key:
            current_round = "round2"
            continue
        if upper == f_key:
            # Final slide content (category can be pulled from title or fixed)
            q = get_body_text(slide)
            a = get_notes_text(slide)

            imgs = extract_images_from_slide(slide, images_dir, idx)
            auds = extract_audio_from_slide(pptx_path, audio_dir, idx)

            final = {
                "category": args.final_title,   # keep explicit
                "question": q,
                "answer": a
            }
            if imgs:
                final["images"] = [{"src": f"{assets_root}/images/{n}"} for n in imgs]
            if auds:
                final["audio"] = [{"src": f"{assets_root}/audio/{n}"} for n in auds]

            current_round = None
            continue

        # Only process clue slides when inside round sections
        if current_round not in ("round1", "round2"):
            continue

        category = title
        question = get_body_text(slide)
        answer = get_notes_text(slide)

        # Skip empty or template/instruction slides
        if not category or not question or not answer:
            continue

        imgs = extract_images_from_slide(slide, images_dir, idx)
        auds = extract_audio_from_slide(pptx_path, audio_dir, idx)

        clue = {
            "question": question,
            "answer": answer
        }
        if imgs:
            clue["images"] = [{"src": f"{assets_root}/images/{n}"} for n in imgs]
        if auds:
            clue["audio"] = [{"src": f"{assets_root}/audio/{n}"} for n in auds]

        if category not in rounds[current_round]:
            rounds[current_round][category] = []
        rounds[current_round][category].append(clue)

    # Build JSON output (content only)
    data = {
        "title": "Jeopardy Content",
        "rounds": [
            {
                "name": "Round 1",
                "categories": [{"name": cat, "clues": clues} for cat, clues in rounds["round1"].items()]
            },
            {
                "name": "Round 2",
                "categories": [{"name": cat, "clues": clues} for cat, clues in rounds["round2"].items()]
            }
        ],
        "final": final
    }

    with open(out_json, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

    print(f"✅ Wrote: {out_json}")
    print(f"✅ Images: {images_dir}")
    print(f"✅ Audio:  {audio_dir}")
    print(f"Round 1 categories: {len(rounds['round1'])}")
    print(f"Round 2 categories: {len(rounds['round2'])}")
    print(f"Final present: {'yes' if final else 'no'}")

if __name__ == "__main__":
    main()